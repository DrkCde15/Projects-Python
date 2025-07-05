# J.A.R.V.I.S = Just A Rather Very Intelligent System

import google.generativeai as genai
import pyttsx3
import webbrowser
import os
import json
import datetime
import re
import speech_recognition as sr
import traceback
import urllib.parse
from dotenv import load_dotenv
import getpass
import hashlib
import subprocess
from pathlib import Path
import fitz
from docx import Document
import pandas as pd
from pptx import Presentation

# ========== CONFIG ==========
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

idioma = 'pt'
voz_ativa = True
COMANDOS_PATH = "./arquivos/sites.json"
APLICATIVOS_PATH = "./arquivos/aplicativos.json"
USUARIOS_ADMIN_PATH = "./arquivos/usuarios_admin.json"

# ========== VOZ ==========
engine = pyttsx3.init()
voices = engine.getProperty('voices')
engine.setProperty('rate', 250)

def configurar_voz():
    for voice in voices:
        if idioma == 'pt' and 'brazil' in voice.name.lower():
            engine.setProperty('voice', voice.id)
            return

def falar(texto):
    if not voz_ativa:
        return
    configurar_voz()
    engine.say(texto)
    engine.runAndWait()

# ========== JSON ==========
def carregar_json(path):
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def salvar_json(path, dados):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(dados, f, indent=4)

comandos_personalizados = carregar_json(COMANDOS_PATH)
aplicativos = carregar_json(APLICATIVOS_PATH)

# ========== GEMINI ==========
def responder_com_gemini(prompt):
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        resposta = model.generate_content(prompt)
        texto = resposta.text.strip()
        print(f"JARVIS: {texto}")
        falar(texto)
        return texto
    except Exception as e:
        return f"Erro Gemini: {e}"

# ========== ADMIN USUÁRIOS ==========
def carregar_usuarios_admin():
    if os.path.exists(USUARIOS_ADMIN_PATH):
        with open(USUARIOS_ADMIN_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def salvar_usuarios_admin(dados):
    with open(USUARIOS_ADMIN_PATH, "w", encoding="utf-8") as f:
        json.dump(dados, f, indent=4)

def hash_senha(senha):
    return hashlib.sha256(senha.encode()).hexdigest()

def cadastrar_usuario_admin():
    usuarios = carregar_usuarios_admin()
    nome = input("Novo nome de usuário admin: ").strip().lower()
    if nome in usuarios:
        print("Usuário já existe.")
        return
    senha = getpass.getpass("Crie uma senha: ").strip()
    senha2 = getpass.getpass("Confirme a senha: ").strip()
    if senha != senha2:
        print("As senhas não coincidem.")
        return
    usuarios[nome] = hash_senha(senha)
    salvar_usuarios_admin(usuarios)
    print(f"Usuário admin '{nome}' cadastrado com sucesso.")

# ========== COMANDOS ==========
urls = {
    'youtube': 'https://youtube.com',
    'netflix': 'https://www.netflix.com/browse',
    'microsoft teams': 'https://teams.microsoft.com/v2/',
    'github': 'https://github.com/DrkCde15',
    'instagram': 'https://www.instagram.com/jc_v05/',
    'whatsapp': 'https://web.whatsapp.com/',
    'tik tok': 'https://www.tiktok.com/@bx_329',
    'e-mail': 'https://mail.google.com/mail/u/1/#inbox'
}

def abrir_site(site):
    if site in urls:
        try:
            webbrowser.open(urls[site])
            return f"Abrindo {site}."
        except Exception as e:
            return f"Erro ao abrir {site}: {str(e)}"
    return "Site não reconhecido."

def abrir_aplicativo(nome):
    nome = nome.lower().strip().split()[0]
    if nome in aplicativos:
        try:
            os.system(aplicativos[nome])
            return f"Abrindo {nome}..."
        except Exception as e:
            return f"Erro ao abrir {nome}: {str(e)}"
    return f"Aplicativo '{nome}' não está mapeado."

def abrir_pasta(nome):
    home = Path.home()
    nome = nome.lower()
    sinonimos = {
        "documentos": ["documento", "documentos", "meus documentos"],
        "imagens": ["imagem", "imagens", "minhas imagens", "fotos"],
        "downloads": ["download", "downloads", "baixados", "transferências"],
        "projetos": ["projeto", "projetos", "meus projetos"],
        "aniversarios": ["aniversario", "aniversarios"],
        "codigos": ["codigo", "codigos", "meus codigos", "scripts", "projetos de código"],
        "desktop": ["desktop", "área de trabalho", "tela inicial"],
        "videos": ["video", "videos", "filmes"],
        "musica": ["musica", "músicas", "audios", "sons"]
    }
    caminhos = {
        "documentos": home / "Documents",
        "imagens": home / "Pictures",
        "downloads": home / "Downloads",
        "projetos": home / "Documents" / "Projects",
        "aniversarios": home / "Documents" / "aniversarios",
        "codigos": home / "Documents" / "Codes-master",
        "desktop": home / "Desktop",
        "videos": home / "Videos",
        "musica": home / "Music"
    }
    pasta_key = None
    for chave, lista in sinonimos.items():
        if nome in lista:
            pasta_key = chave
            break
    if not pasta_key:
        return f"Pasta '{nome}' não reconhecida."
    caminho = caminhos.get(pasta_key)
    if not caminho.exists():
        return f"Pasta '{pasta_key}' mapeada, mas o caminho '{caminho}' não existe."
    try:
        subprocess.Popen(f'explorer "{caminho}"', shell=True)
        return f"Abrindo pasta {pasta_key}."
    except Exception as e:
        return f"Erro ao tentar abrir a pasta '{pasta_key}': {e}"

def falar_hora():
    return f"Agora são {datetime.datetime.now().strftime('%H:%M')}."

def falar_data():
    return f"Hoje é {datetime.datetime.now().strftime('%d/%m/%Y')}."

def listar_aplicativos():
    if not aplicativos:
        return "Nenhum aplicativo foi mapeado ainda."
    return "Aplicativos disponíveis: " + ", ".join(sorted(aplicativos.keys()))

def listar_sites():
    if not urls:
        return "Nenhum site foi mapeado ainda."
    return "Sites disponíveis: " + ", ".join(sorted(urls.keys()))

def pesquisar_google(termo):
    termo_codificado = urllib.parse.quote_plus(termo)
    url = f"https://www.google.com/search?q={termo_codificado}"
    webbrowser.open(url)
    return f"Pesquisando '{termo}'"

def ler_docx(caminho):
    try:
        doc = Document(caminho)
        texto = '\n'.join([p.text for p in doc.paragraphs])
        return texto
    except Exception as e:
        return f"Erro ao ler .docx: {e}"

def ler_pdf(caminho):
    try:
        doc = fitz.open(caminho)
        texto = ''
        for pagina in doc:
            texto += pagina.get_text()
        return texto
    except Exception as e:
        return f"Erro ao ler PDF: {e}"

def ler_pdf(caminho):
    try:
        doc = fitz.open(caminho)
        texto = ''
        for pagina in doc:
            texto += pagina.get_text()
        return texto
    except Exception as e:
        return f"Erro ao ler PDF: {e}"

def ler_excel(caminho):
    try:
        df = pd.read_excel(caminho)
        return df.to_string(index=False)
    except Exception as e:
        return f"Erro ao ler Excel: {e}"

def ler_pptx(caminho):
    try:
        prs = Presentation(caminho)
        texto = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto
    except Exception as e:
        return f"Erro ao ler PPTX: {e}"

def buscar_arquivos_recursivo(pasta, padrao='*'):
    p = Path(pasta)
    arquivos = list(p.rglob(padrao))
    return arquivos

def listar_arquivos_extensao(pasta, extensao):
    padrao = f'*.{extensao.lstrip(".")}'
    arquivos = buscar_arquivos_recursivo(pasta, padrao)
    if not arquivos:
        return f"Nenhum arquivo com extensão {extensao} encontrado em {pasta}."
    arquivos_str = "\n".join(str(a) for a in arquivos[:20])
    total = len(arquivos)
    mais = f"\n...e mais {total - 20} arquivos." if total > 20 else ""
    return f"Arquivos encontrados ({total}):\n{arquivos_str}{mais}"

def analisar_arquivos(comando):
    if comando.startswith("analisar arquivo"):
        try:
            caminho = comando.replace("analisar arquivo", "").strip('" ').strip()
            if not os.path.exists(caminho):
                return "Arquivo não encontrado."
            if caminho.endswith(".pdf"):
                conteudo = ler_pdf(caminho)
            elif caminho.endswith(".docx"):
                conteudo = ler_docx(caminho)
            elif caminho.endswith(".pptx"):
                conteudo = ler_pptx(caminho)
            elif caminho.endswith(".xlsx") or caminho.endswith(".xls"):
                conteudo = ler_excel(caminho)
            else:
                return "Formato ainda não suportado para análise."

            prompt = f"Analise esse conteúdo extraído do arquivo:\n\n{conteudo}"
            return responder_com_gemini(prompt)
        except Exception as e:
            return f"Erro ao analisar arquivo: {e}"
    return None

# ========== PADRÕES ==========
padroes = [
    (r'\b(iniciar|abrir|executar)\s+(youtube|netflix|microsoft teams|github|instagram|whatsapp|tik tok|e-mail)', lambda m: abrir_site(m.group(2))),
    (r'\b(executar|abrir|iniciar)\s+([a-zA-Z0-9_ ]+)', lambda m: abrir_aplicativo(m.group(2))),
    (r'\b(abrir|acessar|mostrar|ver|exibir|quero abrir|abre|abrir a|abrir os|abrir as|mostrar os|mostrar as|acessar os|acessar as)\s+(a|o|os|as|meu|meus|minha|minhas)?\s*([\w\s]+)', lambda m: abrir_pasta(m.group(3).strip())),
    (r'\b(analisar|interpretar|ler|estudar)\s+arquivo\s+(.+)', lambda m: analisar_arquivos(f'analisar arquivo {m.group(2).strip()}')),
    (r'\b(que horas( são)?|horas|hora atual|me diga as horas|qual é a hora)\b', lambda m: falar_hora()),
    (r'\b(data|que dia é hoje|me diga a data|qual a data|qual é a data)\b', lambda m: falar_data()),
    (r'\b(listar|quais são os|me diga os|mostrar)\s+(aplicativo|apps)\b', lambda m: listar_aplicativos()),
    (r'\b(listar|quais são os|me diga os|mostrar)\s+(site|endereços|links)\b', lambda m: listar_sites()),
    (r'pesquisar por\s+(.+)', lambda m: pesquisar_google(m.group(1).strip()))
]

def processar_regex(comando):
    for padrao, acao in padroes:
        match = re.search(padrao, comando, re.IGNORECASE)
        if match:
            return acao(match)
    return None

def executar_comando(comando):
    comando = comando.lower().strip()

    if comando in ['abrir e-mail', 'abrir email', 'acessar e-mail']:
        try:
            tipo_msg = input("Que tipo de mensagem você deseja enviar? (ex: profissional para o chefe, agradecimento, etc): ").strip()
            if not tipo_msg:
                return "Tipo de mensagem não especificado."
            prompt = f"Crie uma mensagem de e-mail em português no estilo: {tipo_msg}."
            mensagem = responder_com_gemini(prompt)
            print("\nModelo sugerido:\n")
            print(mensagem)
            falar("Mensagem gerada. Abrindo Gmail.")
            webbrowser.open(urls['e-mail'])
            return "Email aberto no navegador."
        except Exception as e:
            return f"Erro ao processar mensagem de e-mail: {str(e)}"

    if comando in ['sair', 'exit']:
        return "Encerrando."

    if comando == "cadastrar admin":
        cadastrar_usuario_admin()
        return "Cadastro de administrador finalizado."

    if comando == "modo administrador":
        falar("Modo administrador solicitado. Por favor, autentique-se.")
        usuarios = carregar_usuarios_admin()
        usuario = input("Usuário: ").strip().lower()
        senha = getpass.getpass("Senha: ").strip()
        senha_hash = hash_senha(senha)
        if usuarios.get(usuario) != senha_hash:
            return "Acesso negado. Usuário ou senha incorretos."
        try:
            caminho_admin = os.path.abspath("admin_actions.py")
            proc = subprocess.Popen(['cmd.exe', '/k', 'python', caminho_admin, usuario])
            proc.wait()
            return "Modo administrador encerrado."
        except Exception as e:
            return f"Erro ao tentar ativar modo administrador: {e}"

    if comando.startswith("cadastrar comando"):
        try:
            partes = comando.replace("cadastrar comando", "").strip().split(" para ")
            nome, acao = partes[0].strip(), partes[1].strip()
            comandos_personalizados[nome] = acao
            salvar_json(COMANDOS_PATH, comandos_personalizados)
            return f"Comando '{nome}' cadastrado com sucesso."
        except:
            return "Formato inválido. Use: cadastrar comando NOME para LINK ou COMANDO."

    if comando in comandos_personalizados:
        destino = comandos_personalizados[comando]
        try:
            if destino.startswith("http"):
                webbrowser.open(destino)
            elif destino.endswith("()") or destino.startswith("abrir_aplicativo"):
                return eval(destino)
            else:
                os.system(destino)
            return f"Executando {comando}."
        except Exception as e:
            return f"Erro ao executar comando: {str(e)}"

    if comando.startswith("listar arquivos"):
        try:
            match = re.search(r"listar arquivos em (\w+)\s+com\s+(?:a\s+)?extens[aã]o\s+(\w+)", comando)
            if not match:
                return "Formato inválido. Use: listar arquivos em [pasta] com a extensão [extensão]"
            pasta = match.group(1)
            extensao = match.group(2)
            home = Path.home()
            pastas_map = {
                "documento": home / "Documents",
                "documentos": home / "Documents",
                "imagem": home / "Pictures",
                "imagens": home / "Pictures",
                "download": home / "Downloads",
                "downloads": home / "Downloads",
                "projeto": home / "Documents" / "Projects",
                "projetos": home / "Documents" / "Projects",
                "aniversario": home / "Documents" / "aniversarios",
                "aniversarios": home / "Documents" / "aniversarios",
                "codigo": home / "Documents" / "Codes-master",
                "codigos": home / "Documents" / "Codes-master",
                "desktop": home / "Desktop",
                "area de trabalho": home / "Desktop"
            }
            caminho = pastas_map.get(pasta, pasta)
            return listar_arquivos_extensao(caminho, extensao)
        except Exception:
            return "Erro ao tentar listar arquivos."

    resposta_regex = processar_regex(comando)
    if resposta_regex:
        return resposta_regex

    return responder_com_gemini(comando)

def ouvir_comando():
    r = sr.Recognizer()
    with sr.Microphone() as source:
        r.adjust_for_ambient_noise(source)
        print("Escutando...")
        try:
            audio = r.listen(source, timeout=5, phrase_time_limit=15)
        except sr.WaitTimeoutError:
            return ""
    try:
        comando = r.recognize_google(audio, language='pt-BR')
        print("Você disse:", comando)
        return comando.lower()
    except:
        return ""

def modo_voz_manual():
    falar("Olá Senhor, Sou seu assistente de voz JARVIS, por favor me diga um comando.")
    while True:
        comando = ouvir_comando()
        if comando:
            if comando in ['sair', 'exit']:
                falar("Até mais Senhor.")
                break
            resposta = executar_comando(comando)
            if resposta:
                falar(resposta)

def modo_texto_terminal():
    print("Ola Senhor, Sou seu assistente JARVIS. Digite 'x' ou 'exit' para encerrar.\n")
    try:
        while True:
            pergunta = input("Usuário: ").strip()
            if pergunta.lower() in ['x', 'exit']:
                print("Até mais Senhor!")
                break
            resposta = executar_comando(pergunta)
            if resposta:
                print(f"\nJARVIS: {resposta}\n")
    except KeyboardInterrupt:
        print("\nInterrupção detectada. Até mais Senhor.")
    except Exception:
        print("Erro inesperado:")
        traceback.print_exc()
        
# Adicione esta função no final do seu script, **antes** do `if __name__ == "__main__":`
def mostrar_manual():
    print("\n========== MANUAL DE USO - J.A.R.V.I.S ==========\n")
    print("JARVIS (Just A Rather Very Intelligent System) é seu assistente pessoal automatizado.")
    print("Use comandos de texto ou voz para interagir com o sistema.")
    print("\n MODO TEXTO:")
    print("- Você digita os comandos diretamente no terminal.")
    print("- Exemplos:")
    print("  → abrir navegador")
    print("  → listar aplicativos")
    print("  → abrir pasta downloads")
    print("  → listar arquivos em documentos com extensão pdf/docx/pptx/xlsx/xls")
    print("  → analisar documentos em pdf/docx/pptx/xlsx/xls")
    print("  → rodar scanner de IP (modo admin)")
    print("  → desligar o sistema")
    
    print("\n MODO VOZ:")
    print("- O JARVIS escuta e interpreta sua fala.")
    print("- Fale um comando por vez com clareza.")
    print("- Exemplos:")
    print("  → abrir navegador")
    print("  → listar aplicativos")
    print("  → abrir pasta downloads")
    
    print("\n MODO ADMIN (restrito):")
    print("- Ativado por usario e senha admin cadastrados.")
    print("- Permite:")
    print("  → Criar usuários (net user)")
    print("  → Alterar registros do sistema")
    print("  → Executar comandos críticos como shutdown, nmap, reg add, winget")
    
    print("\n COMANDOS DISPONÍVEIS:")
    print("- abrir navegador / email / whatsapp / youtube")
    print("- abrir pasta [nome]")
    print("- listar sites / aplicativos")
    print("- checar atualizações do Windows")
    print("- abrir aplicativo pelo nome")
    print("- analisar arquivos PDF, DOCX, PPTX e planilhas Excel")
    
    print("\n DICAS:")
    print("- Fale de forma objetiva no modo de voz.")
    print("- Evite ruídos ou interrupções no microfone.")
    print("- Comandos complexos exigem o modo admin ativado.")
    
    print("\n=================================================\n")
    
if __name__ == "__main__":
    print("\nEscolha a forma de interação:")
    print("1 - Modo por voz (fala um comando por vez)")
    print("2 - Modo texto (JARVIS)")
    print("3 - Ver manual de uso do JARVIS")

    escolha = input("Digite 1, 2 ou 3: ").strip()

    if escolha == '1':
        voz_ativa = True
        modo_voz_manual()
    elif escolha == '2':
        voz_ativa = False
        modo_texto_terminal()
    elif escolha == '3':
        mostrar_manual()
    else:
        print("Opção inválida.")